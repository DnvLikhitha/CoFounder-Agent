"""
export/pptx_generator.py — Generate a .pptx pitch deck using python-pptx
"""

from __future__ import annotations

from io import BytesIO

from pptx import Presentation
from pptx.util import Inches, Pt

from backend.context import RunContext


def _add_title(slide, title: str) -> None:
    """Add a title text box at the top of the slide."""
    if not title:
        title = "Pitch Deck"
    tx = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12.0), Inches(0.7))
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(28)


def _add_subtitle(slide, subtitle: str) -> None:
    """Add subtitle text under the title."""
    if not subtitle:
        return
    tx = slide.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(12.0), Inches(0.5))
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = subtitle
    run.font.size = Pt(16)


def _add_bullets(slide, bullets: list[str]) -> None:
    """Add bullet list on the left."""
    if not bullets:
        return
    tx = slide.shapes.add_textbox(Inches(0.9), Inches(2.0), Inches(11.3), Inches(4.8))
    tf = tx.text_frame
    tf.clear()
    tf.word_wrap = True

    # First paragraph created by default
    for i, b in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = b
        p.level = 0
        for run in p.runs:
            run.font.size = Pt(18)


def _add_footer(slide, text: str) -> None:
    """Add a small footer at the bottom."""
    if not text:
        return
    tx = slide.shapes.add_textbox(Inches(0.6), Inches(6.9), Inches(12.0), Inches(0.3))
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.runs[0].font.size = Pt(10)


def generate_pptx(ctx: RunContext) -> bytes:
    """
    Generate a pitch deck PPTX from ctx.pitch_deck.slides (Agent14 output).
    """
    prs = Presentation()

    idea = ctx.startup_idea or {}
    deck = ctx.pitch_deck or {}
    slides = deck.get("slides", [])

    # Sort defensively by slide_number
    def slide_key(s: dict) -> int:
        try:
            return int(s.get("slide_number"))
        except Exception:
            return 999

    slides_sorted = sorted(slides, key=slide_key)

    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    _add_title(title_slide, idea.get("startup_name", "Startup Plan"))
    _add_subtitle(title_slide, idea.get("tagline", "Investor-ready pitch deck"))
    _add_footer(title_slide, f"AI Startup Co-Founder Agent | Run ID: {ctx.run_id}")

    # Remaining slides (skip slides[0] if it represents "The Problem" etc; Agent still includes 12 slides)
    for s in slides_sorted:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
        title = s.get("title") or s.get("headline") or f"Slide {s.get('slide_number', '')}"
        _add_title(slide, title)

        bullets = s.get("bullet_points") or []
        # If Agent provides headline, include it as the first bullet contextually.
        headline = s.get("headline")
        if headline and isinstance(headline, str) and headline.strip():
            bullets = [headline] + bullets

        _add_bullets(slide, bullets[:8])

        key_stat = s.get("key_stat")
        if key_stat:
            _add_subtitle(slide, str(key_stat)[:120])

        _add_footer(slide, f"Run ID: {ctx.run_id}")

    out = BytesIO()
    prs.save(out)
    out.seek(0)
    return out.read()

